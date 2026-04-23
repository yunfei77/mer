from email import policy
from email.parser import BytesParser
import extract_msg
import os
from typing import Dict, Any, List
from bs4 import BeautifulSoup
import re
import io
import time
try:
    import PyPDF2  # 用于解析PDF文件
    from docx import Document  # 用于解析Word文档
    import openpyxl  # 用于解析Excel文件
    import pptx  # 用于解析PPT文件
    EXTRA_FORMATS_SUPPORTED = True
except ImportError:
    EXTRA_FORMATS_SUPPORTED = False
    print("提示: 要支持更多文档格式预览，请安装以下包:")
    print("pip install python-docx PyPDF2 openpyxl python-pptx")

import whois
from datetime import datetime, timezone

def parse_email(file_path: str) -> Dict[str, Any]:
    """
    解析邮件文件，提取关键信息
    
    Args:
        file_path: 邮件文件路径
        
    Returns:
        包含邮件信息的字典
    """
    email_data = {
        'from': [],
        'to': [],
        'cc': [],
        'bcc': [],
        'reply_to': [],
        'subject': '',
        'date': '',
        'body_text': '',
        'body_html': '',
        'attachments': [],
        'references': [],
        'in_reply_to': [],
        'headers': {},        # 新增：完整原始 headers
        'thread_info': {
            'original_sender': '',
            'original_recipients': [],
            'original_subject': '',
            'original_date': ''
        }
    }
    
    try:
        if file_path.lower().endswith('.eml'):
            with open(file_path, 'rb') as f:
                msg = BytesParser(policy=policy.default).parse(f)
                
                # 解析邮件头
                def parse_address_list(header_value):
                    """解析邮件地址列表"""
                    if not header_value:
                        return []
                    # 处理可能的多行地址
                    if isinstance(header_value, (list, tuple)):
                        addresses = []
                        for item in header_value:
                            addresses.extend([addr.strip() for addr in str(item).split(',')])
                        return addresses
                    return [addr.strip() for addr in str(header_value).split(',')]
                
                # 处理发件人
                from_header = msg.get('From', '')
                email_data['from'] = parse_address_list(from_header)
                
                # 处理收件人
                to_header = msg.get('To', '')
                email_data['to'] = parse_address_list(to_header)
                
                # 处理抄送人
                cc_header = msg.get('Cc', '')
                email_data['cc'] = parse_address_list(cc_header)
                
                # 处理密送人
                bcc_header = msg.get('Bcc', '')
                email_data['bcc'] = parse_address_list(bcc_header)
                
                # 处理回复地址
                reply_to_header = msg.get('Reply-To', '')
                email_data['reply_to'] = parse_address_list(reply_to_header)
                
                # 处理主题和日期
                email_data['subject'] = msg.get('Subject', '')
                email_data['date'] = msg.get('Date', '')

                # 提取所有原始 headers（供 SPF/DKIM/DMARC/Received 检测使用）
                for key in msg.keys():
                    k = key.lower()
                    val = str(msg.get(key, ''))
                    if k == 'received':
                        email_data['headers'].setdefault('received', [])
                        if isinstance(email_data['headers']['received'], list):
                            email_data['headers']['received'].append(val)
                    else:
                        email_data['headers'][k] = val
                
                # 处理邮件引用和回复信息
                references = msg.get('References', '')
                in_reply_to = msg.get('In-Reply-To', '')
                email_data['references'] = parse_address_list(references)
                email_data['in_reply_to'] = parse_address_list(in_reply_to)
                
                # 尝试获取原始邮件信息
                if msg.is_multipart():
                    for part in msg.walk():
                        if part.get_content_type() == 'message/rfc822':
                            # 获取原始邮件
                            original_msg = part.get_payload()[0]
                            email_data['thread_info']['original_sender'] = original_msg.get('From', '')
                            email_data['thread_info']['original_recipients'] = parse_address_list(original_msg.get('To', ''))
                            email_data['thread_info']['original_subject'] = original_msg.get('Subject', '')
                            email_data['thread_info']['original_date'] = original_msg.get('Date', '')
                
                # 改进附件解析部分
                def extract_attachments(msg):
                    """递归提取附件"""
                    if msg.is_multipart():
                        for part in msg.walk():
                            if part.get_content_maintype() == 'multipart':
                                continue
                            
                            disposition = part.get_content_disposition()
                            if disposition and disposition.lower() in ['attachment', 'inline']:
                                try:
                                    attachment_info = parse_attachment(part, len(email_data['attachments']))
                                    email_data['attachments'].append(attachment_info)
                                    print(f"发现{'内联' if attachment_info['is_inline'] else ''}附件: {attachment_info['filename']}")
                                except Exception as e:
                                    print(f"处理附件出错: {str(e)}")
                
                # 调用附件提取函数
                extract_attachments(msg)
                
                # 改进正文解析部分
                def get_body_content(msg):
                    """递归获取邮件正文内容"""
                    if msg.is_multipart():
                        # 获取所有部分
                        all_parts = []
                        for part in msg.walk():
                            if part.get_content_maintype() == 'text':
                                all_parts.append(part)
                        
                        # 优先处理 text/plain
                        for part in all_parts:
                            if part.get_content_type() == 'text/plain':
                                try:
                                    payload = part.get_payload(decode=True)
                                    if payload:
                                        # 尝试多种编码方式
                                        encodings = ['utf-8', 'gbk', 'gb2312', 'gb18030', 'big5']
                                        for encoding in encodings:
                                            try:
                                                decoded_text = payload.decode(encoding)
                                                if decoded_text:
                                                    email_data['body_text'] += decoded_text + '\n'
                                                    break
                                            except:
                                                continue
                                except Exception as e:
                                    print(f"解析纯文本内容失败: {str(e)}")
                        
                        # 然后处理 text/html
                        for part in all_parts:
                            if part.get_content_type() == 'text/html':
                                try:
                                    # 检查内容传输编码
                                    transfer_encoding = part.get('Content-Transfer-Encoding', '').lower()
                                    
                                    # 获取原始负载
                                    if transfer_encoding == 'base64':
                                        payload = part.get_payload(decode=True)
                                    else:
                                        payload = part.get_payload(decode=True)
                                    
                                    if payload:
                                        # 尝试多种编码方式
                                        encodings = ['utf-8', 'gbk', 'gb2312', 'gb18030', 'big5']
                                        for encoding in encodings:
                                            try:
                                                decoded_html = payload.decode(encoding)
                                                if decoded_html:
                                                    # 验证是否为有效的HTML内容
                                                    if any(marker in decoded_html.lower() for marker in ['<html', '<!doctype', '<body', '<head', '<div', '<p']):
                                                        email_data['body_html'] = decoded_html
                                                        break
                                            except:
                                                continue
                                except Exception as e:
                                    print(f"解析HTML内容失败: {str(e)}")
                    else:
                        # 非多部分邮件
                        content_type = msg.get_content_type()
                        try:
                            # 检查内容传输编码
                            transfer_encoding = msg.get('Content-Transfer-Encoding', '').lower()
                            
                            # 获取原始负载
                            if transfer_encoding == 'base64':
                                payload = msg.get_payload(decode=True)
                            else:
                                payload = msg.get_payload(decode=True)
                            
                            if payload:
                                # 尝试多种编码方式
                                encodings = ['utf-8', 'gbk', 'gb2312', 'gb18030', 'big5']
                                for encoding in encodings:
                                    try:
                                        decoded_content = payload.decode(encoding)
                                        if decoded_content:
                                            if content_type == 'text/plain':
                                                email_data['body_text'] = decoded_content
                                                break
                                            elif content_type == 'text/html':
                                                if any(marker in decoded_content.lower() for marker in ['<html', '<!doctype', '<body', '<head', '<div', '<p']):
                                                    email_data['body_html'] = decoded_content
                                                    break
                                    except:
                                        continue
                        except Exception as e:
                            print(f"解析正文出错: {str(e)}")
                
                # 调用正文解析函数
                get_body_content(msg)
                
                # 如果正文仍然为空，尝试其他方法
                if not email_data['body_text'] and not email_data['body_html']:
                    try:
                        # 直接获取原始负载
                        payload = msg.get_payload()
                        if isinstance(payload, list):
                            for part in payload:
                                if part.get_content_type() == 'text/plain':
                                    content = part.get_payload(decode=True)
                                    if content:
                                        email_data['body_text'] += content.decode('utf-8', errors='ignore') + '\n'
                                elif part.get_content_type() == 'text/html':
                                    content = part.get_payload(decode=True)
                                    if content:
                                        email_data['body_html'] += content.decode('utf-8', errors='ignore') + '\n'
                        elif isinstance(payload, str):
                            email_data['body_text'] = payload
                    except Exception as e:
                        print(f"备用方法获取正文失败: {str(e)}")
                
                # 如果HTML内容存在但看起来是乱码，尝试不同的解码方式
                if email_data['body_html'] and not email_data['body_html'].strip().startswith(('<html', '<!DOCTYPE', '<body', '<div')):
                    try:
                        # 尝试不同的编码
                        encodings = ['utf-8', 'gbk', 'gb2312', 'gb18030', 'big5', 'latin1']
                        original_content = email_data['body_html'].encode('latin1', errors='ignore')
                        
                        for encoding in encodings:
                            try:
                                decoded = original_content.decode(encoding)
                                if decoded.strip().startswith(('<html', '<!DOCTYPE', '<body', '<div')):
                                    email_data['body_html'] = decoded
                                    break
                            except:
                                continue
                    except Exception as e:
                        print(f"尝试重新解码HTML内容失败: {str(e)}")
                
        elif file_path.lower().endswith('.msg'):
            msg = extract_msg.Message(file_path)
            try:
                # 处理发件人
                sender = msg.sender
                if sender:
                    email_data['from'].append(sender)
                
                # 处理收件人
                recipients = msg.to
                if recipients:
                    email_data['to'] = [r.strip() for r in recipients.split(';')]
                
                # 处理抄送人
                cc = msg.cc
                if cc:
                    email_data['cc'] = [c.strip() for c in cc.split(';')]
                
                # 处理密送人
                try:
                    bcc = msg.bcc
                    if bcc:
                        email_data['bcc'] = [b.strip() for b in bcc.split(';')]
                except AttributeError:
                    print("MSG文件不包含密送人信息")
                
                # 处理主题和日期
                email_data['subject'] = msg.subject or ''
                email_data['date'] = str(msg.date) if msg.date else ''
                
                # 改进原始邮件信息获取
                try:
                    # 从邮件正文中提取原始邮件信息
                    body_text = msg.body or ''
                    
                    # 定义更强大的正则表达式模式
                    from_pattern = r'From:[\s]*([^\r\n]+)'
                    to_pattern = r'To:[\s]*([^\r\n]+)'
                    subject_pattern = r'Subject:[\s]*([^\r\n]+)'
                    sent_pattern = r'Sent:[\s]*([^\r\n]+)'
                    cc_pattern = r'Cc:[\s]*([^\r\n]+)'
                    
                    # 查找所有匹配项（可能有多个，取最后一个作为原始邮件信息）
                    from_matches = re.findall(from_pattern, body_text)
                    to_matches = re.findall(to_pattern, body_text)
                    subject_matches = re.findall(subject_pattern, body_text)
                    sent_matches = re.findall(sent_pattern, body_text)
                    cc_matches = re.findall(cc_pattern, body_text)
                    
                    # 如果找到匹配项，使用最后一个（通常是原始邮件的信息）
                    if from_matches:
                        email_data['thread_info']['original_sender'] = from_matches[-1].strip()
                    
                    if to_matches:
                        recipients = [r.strip() for r in to_matches[-1].split(';')]
                        email_data['thread_info']['original_recipients'] = recipients
                    
                    if subject_matches:
                        email_data['thread_info']['original_subject'] = subject_matches[-1].strip()
                    
                    if sent_matches:
                        email_data['thread_info']['original_date'] = sent_matches[-1].strip()
                    
                    if cc_matches:
                        cc_list = [c.strip() for c in cc_matches[-1].split(';')]
                        if not email_data['cc']:  # 如果之前没有抄送人，则添加
                            email_data['cc'] = cc_list
                        else:  # 如果已有抄送人，则合并并去重
                            email_data['cc'].extend(cc_list)
                            email_data['cc'] = list(set(email_data['cc']))
                    
                    # 如果在正文中没有找到，尝试从HTML正文中提取
                    if not any(email_data['thread_info'].values()):
                        html_body = msg.htmlBody or ''
                        if html_body:
                            # 使用BeautifulSoup清理HTML标签
                            soup = BeautifulSoup(html_body, 'html.parser')
                            clean_text = soup.get_text()
                            
                            # 重新查找
                            from_matches = re.findall(from_pattern, clean_text)
                            to_matches = re.findall(to_pattern, clean_text)
                            subject_matches = re.findall(subject_pattern, clean_text)
                            sent_matches = re.findall(sent_pattern, clean_text)
                            
                            if from_matches:
                                email_data['thread_info']['original_sender'] = from_matches[-1].strip()
                            if to_matches:
                                recipients = [r.strip() for r in to_matches[-1].split(';')]
                                email_data['thread_info']['original_recipients'] = recipients
                            if subject_matches:
                                email_data['thread_info']['original_subject'] = subject_matches[-1].strip()
                            if sent_matches:
                                email_data['thread_info']['original_date'] = sent_matches[-1].strip()
                
                except Exception as e:
                    print(f"获取原始邮件信息失败: {str(e)}")
                
                # 处理正文
                email_data['body_text'] = msg.body or ''
                email_data['body_html'] = msg.htmlBody or ''
                
                # 处理附件
                for attachment in msg.attachments:
                    try:
                        attachment_info = parse_attachment(attachment, len(email_data['attachments']))
                        email_data['attachments'].append(attachment_info)
                        print(f"发现MSG附件: {attachment_info['filename']}")
                    except Exception as e:
                        print(f"处理MSG附件出错: {str(e)}")
                
                msg.close()
                
            except Exception as e:
                print(f"处理MSG文件出错: {str(e)}")
                if msg:
                    msg.close()
                
    except Exception as e:
        print(f"邮件解析失败: {str(e)}")
        
    return email_data

def parse_attachment(attachment: Any, attachment_index: int) -> Dict[str, Any]:
    """
    解析邮件附件，提取详细信息
    
    Args:
        attachment: 附件对象
        attachment_index: 附件索引号
        
    Returns:
        包含附件详细信息的字典
    """
    attachment_info = {
        'filename': '',          # 文件名
        'mime_type': '',         # MIME类型
        'size': 0,              # 文件大小
        'data': None,           # 原始数据
        'is_inline': False,     # 是否为内联附件
        'content_id': '',       # 内联附件的Content-ID
        'extension': '',        # 文件扩展名
        'hash_md5': '',         # MD5哈希值
        'hash_sha256': '',      # SHA256哈希值
        'text_preview': '',     # 文本预览（如果是文本文件）
        'is_archive': False,    # 是否为压缩文件
        'archive_contents': [],  # 压缩文件内容列表
        'is_executable': False, # 是否为可执行文件
        'created_date': '',     # 创建日期
        'modified_date': '',    # 修改日期
    }
    
    try:
        # 处理MSG附件
        if hasattr(attachment, 'longFilename'):
            filename = attachment.longFilename or attachment.shortFilename or f"未命名附件_{attachment_index}"
            attachment_info['filename'] = filename
            attachment_info['mime_type'] = attachment.mimetype or 'application/octet-stream'
            attachment_info['data'] = attachment.data
            attachment_info['size'] = len(attachment.data) if attachment.data else 0
            
        # 处理EML附件
        else:
            # 获取文件名
            filename = attachment.get_filename()
            if not filename:
                filename = attachment.get_param('name')
            if not filename:
                filename = f"未命名附件_{attachment_index}"
            attachment_info['filename'] = filename
            
            # 获取MIME类型
            attachment_info['mime_type'] = attachment.get_content_type()
            
            # 获取是否为内联附件
            disposition = attachment.get_content_disposition()
            attachment_info['is_inline'] = (disposition and disposition.lower() == 'inline')
            
            # 获取Content-ID
            content_id = attachment.get('Content-ID', '')
            if content_id:
                attachment_info['content_id'] = content_id.strip('<>')
            
            # 获取附件数据
            payload = attachment.get_payload(decode=True)
            if payload:
                attachment_info['data'] = payload
                attachment_info['size'] = len(payload)
        
        # 获取文件扩展名
        attachment_info['extension'] = os.path.splitext(attachment_info['filename'])[1].lower()
        
        # 计算哈希值
        if attachment_info['data']:
            import hashlib
            attachment_info['hash_md5'] = hashlib.md5(attachment_info['data']).hexdigest()
            attachment_info['hash_sha256'] = hashlib.sha256(attachment_info['data']).hexdigest()
        
        # 检查是否为可执行文件
        executable_extensions = {'.exe', '.dll', '.bat', '.cmd', '.msi', '.vbs', '.js', '.ps1', '.com', '.scr'}
        attachment_info['is_executable'] = attachment_info['extension'] in executable_extensions
        
        # 检查是否为压缩文件并提取内容列表
        archive_extensions = {'.zip', '.rar', '.7z', '.tar', '.gz', '.bz2'}
        if attachment_info['extension'] in archive_extensions:
            attachment_info['is_archive'] = True
            try:
                import zipfile
                if attachment_info['extension'] == '.zip' and attachment_info['data']:
                    from io import BytesIO
                    with zipfile.ZipFile(BytesIO(attachment_info['data'])) as zf:
                        attachment_info['archive_contents'] = zf.namelist()
            except Exception as e:
                print(f"解析压缩文件内容失败: {str(e)}")
        
        # 扩展文本文件和文档预览支持
        preview_extensions = {
            'text': {'.txt', '.csv', '.log', '.xml', '.json', '.html', '.htm', '.css', '.js', '.py', '.java', '.cpp', '.c', '.h', '.sql'},
            'document': {'.pdf', '.docx', '.xlsx', '.pptx', '.doc', '.xls', '.ppt'},
        }
        
        # 如果是文本文件，添加预览
        if attachment_info['extension'] in preview_extensions['text'] and attachment_info['data']:
            try:
                preview_text = attachment_info['data'].decode('utf-8', errors='ignore')
                attachment_info['text_preview'] = preview_text[:2000] + '...' if len(preview_text) > 2000 else preview_text
            except Exception as e:
                print(f"生成文本预览失败: {str(e)}")
        
        # 如果是文档文件且支持扩展格式，添加预览
        elif EXTRA_FORMATS_SUPPORTED and attachment_info['extension'] in preview_extensions['document'] and attachment_info['data']:
            try:
                preview_text = ""
                
                # PDF文件处理
                if attachment_info['extension'] == '.pdf':
                    try:
                        pdf_file = io.BytesIO(attachment_info['data'])
                        pdf_reader = PyPDF2.PdfReader(pdf_file)
                        preview_text = "PDF文档内容预览:\n"
                        # 只预览前3页
                        for page_num in range(min(3, len(pdf_reader.pages))):
                            preview_text += f"\n--- 第{page_num + 1}页 ---\n"
                            preview_text += pdf_reader.pages[page_num].extract_text()[:1000]
                            if page_num < min(2, len(pdf_reader.pages) - 1):
                                preview_text += "\n...\n"
                    except Exception as e:
                        print(f"PDF文档预览失败: {str(e)}")
                
                # Word文档处理
                elif attachment_info['extension'] == '.docx':
                    try:
                        docx_file = io.BytesIO(attachment_info['data'])
                        doc = Document(docx_file)
                        preview_text = "Word文档内容预览:\n\n"
                        # 获取文档的前10个段落
                        for i, para in enumerate(doc.paragraphs[:10]):
                            if para.text.strip():
                                preview_text += para.text + "\n"
                        if len(doc.paragraphs) > 10:
                            preview_text += "\n... (更多内容已省略)"
                    except Exception as e:
                        print(f"Word文档预览失败: {str(e)}")
                
                # Excel文件处理
                elif attachment_info['extension'] == '.xlsx':
                    try:
                        xlsx_file = io.BytesIO(attachment_info['data'])
                        wb = openpyxl.load_workbook(xlsx_file, read_only=True)
                        preview_text = "Excel文档内容预览:\n\n"
                        # 预览第一个工作表的前10行
                        sheet = wb.active
                        for i, row in enumerate(sheet.iter_rows(max_row=10)):
                            if i == 0:
                                preview_text += "表头: "
                            preview_text += " | ".join(str(cell.value) for cell in row) + "\n"
                        if sheet.max_row > 10:
                            preview_text += "\n... (更多行已省略)"
                        wb.close()  # 确保关闭工作簿
                    except Exception as e:
                        print(f"Excel文档预览失败: {str(e)}")
                
                # PowerPoint文件处理
                elif attachment_info['extension'] == '.pptx':
                    try:
                        pptx_file = io.BytesIO(attachment_info['data'])
                        prs = pptx.Presentation(pptx_file)
                        preview_text = "PowerPoint文档内容预览:\n\n"
                        # 预览前3张幻灯片
                        for i, slide in enumerate(prs.slides[:3]):
                            preview_text += f"\n--- 幻灯片 {i+1} ---\n"
                            for shape in slide.shapes:
                                if hasattr(shape, "text"):
                                    preview_text += shape.text + "\n"
                        if len(prs.slides) > 3:
                            preview_text += "\n... (更多幻灯片已省略)"
                    except Exception as e:
                        print(f"PowerPoint文档预览失败: {str(e)}")
                
                if preview_text:
                    attachment_info['text_preview'] = preview_text
                
            except Exception as e:
                print(f"文档预览处理失败: {str(e)}")
        
    except Exception as e:
        print(f"解析附件失败: {str(e)}")
    
    return attachment_info

# ========== 公共工具函数 ==========

def levenshtein_distance(s1: str, s2: str) -> int:
    """计算两个字符串的编辑距离"""
    if len(s1) < len(s2):
        return levenshtein_distance(s2, s1)
    if len(s2) == 0:
        return len(s1)
    previous_row = range(len(s2) + 1)
    for i, c1 in enumerate(s1):
        current_row = [i + 1]
        for j, c2 in enumerate(s2):
            current_row.append(min(
                previous_row[j + 1] + 1,
                current_row[j] + 1,
                previous_row[j] + (c1 != c2)
            ))
        previous_row = current_row
    return previous_row[-1]


def analyze_domain_similarity(domain1: str, domain2: str) -> Dict[str, Any]:
    """
    分析两个域名之间的相似关系。
    返回 {'similarity': float, 'relationship_type': str|None, 'risk': str}
    """
    if not domain1 or not domain2 or domain1 == domain2:
        return {'similarity': 0.0, 'relationship_type': None, 'risk': 'low'}

    base1 = domain1.split('.')[0].lower()
    base2 = domain2.split('.')[0].lower()

    # 1. 字符替换（sinopc → sinopec，差1个字符）
    if abs(len(base1) - len(base2)) <= 1:
        diff = sum(c1 != c2 for c1, c2 in zip(base1, base2))
        if diff <= 1:
            return {'similarity': 0.95, 'relationship_type': '可疑的字符替换', 'risk': 'high'}

    # 2. 字母顺序调换（anagram）
    if sorted(base1) == sorted(base2) and base1 != base2:
        return {'similarity': 1.0, 'relationship_type': '字母顺序调换', 'risk': 'high'}

    # 3. 包含关系 + 可疑追加词
    suspicious_additions = {
        'portal', 'service', 'vendor', 'secure', 'mail',
        'auth', 'login', 'account', 'verify', 'update', 'support'
    }
    if base1 in base2 or base2 in base1:
        longer  = base1 if len(base1) > len(base2) else base2
        shorter = base2 if len(base1) > len(base2) else base1
        remaining = longer.replace(shorter, '').lower()
        if any(w in remaining for w in suspicious_additions):
            return {'similarity': 0.9, 'relationship_type': '可疑的域名包含', 'risk': 'high'}

    # 4. 编辑距离相似度
    dist = levenshtein_distance(base1, base2)
    max_len = max(len(base1), len(base2))
    similarity = 1 - dist / max_len if max_len else 0
    if similarity > 0.8:
        return {'similarity': similarity, 'relationship_type': '高度相似', 'risk': 'high'}

    return {'similarity': similarity, 'relationship_type': None, 'risk': 'low'}


def extract_email_domain(email_str: str) -> str:
    """从邮箱地址字符串中提取域名"""
    m = re.search(r'@([\w.-]+)', email_str)
    return m.group(1).lower() if m else ''


def extract_url_domain(url: str) -> str:
    """从 URL 中提取域名"""
    try:
        from urllib.parse import urlparse
        return urlparse(url).netloc.lower()
    except Exception:
        return ''


def risk_label(score: float) -> str:
    """根据分数返回风险等级标签"""
    if score >= 7:   return 'critical'
    if score >= 5:   return 'high'
    if score >= 2.5: return 'medium'
    if score > 0:    return 'low'
    return 'low'


def fmt_risk(level: str) -> str:
    """带颜色的风险等级字符串（ANSI）"""
    colors = {
        'critical': '\033[1;35m',  # 紫
        'high':     '\033[1;31m',  # 红
        'medium':   '\033[1;33m',  # 黄
        'low':      '\033[1;32m',  # 绿
        'unknown':  '\033[0m',
    }
    reset = '\033[0m'
    c = colors.get(level.lower(), reset)
    return f"{c}{level.upper()}{reset}"


def _sep(title: str = '', width: int = 60) -> str:
    if title:
        pad = (width - len(title) - 2) // 2
        return '─' * pad + f' {title} ' + '─' * pad
    return '─' * width


# ========== 新增检测方法 ==========

def detect_homograph_attack(email_data: Dict[str, Any]) -> Dict[str, Any]:
    """
    检测同形字攻击（Homograph/IDN Homograph Attack）：
    使用 Unicode 字符（如西里尔字母 а≠a）伪装域名。
    """
    result = {
        'suspicious': [],
        'risk_level': 'low',
        'risk_score': 0.0,
        'warnings': []
    }

    # ASCII 范围之外的字母映射（常见混淆字符集）
    confusables = {
        'а': 'a', 'е': 'e', 'о': 'o', 'р': 'p', 'с': 'c',
        'х': 'x', 'у': 'y', 'і': 'i', 'ԁ': 'd', 'ɡ': 'g',
        'ⅼ': 'l', '０': '0', '１': '1', '２': '2',
    }

    def has_confusable(s: str) -> bool:
        return any(c in confusables for c in s)

    addresses = (
        email_data.get('from', []) +
        email_data.get('to', []) +
        email_data.get('reply_to', [])
    )
    for addr in addresses:
        if has_confusable(addr):
            result['suspicious'].append(addr)
            result['risk_score'] += 4.0
            result['warnings'].append(
                f'地址含同形混淆字符（可能是仿冒域名）: {addr}'
            )

    if result['risk_score'] >= 4:
        result['risk_level'] = 'high'

    return result


def detect_suspicious_subject(email_data: Dict[str, Any]) -> Dict[str, Any]:
    """
    基于主题关键词的威胁评分。
    涵盖：紧迫感、金融诱导、账户威胁、凭证钓鱼等常见主题词。
    """
    result = {
        'matched_keywords': [],
        'risk_level': 'low',
        'risk_score': 0.0,
        'warnings': []
    }

    HIGH_RISK = [
        '紧急', '立即', '马上', 'urgent', 'immediate', 'action required',
        '账号被冻结', '账户异常', '密码过期', 'password expired',
        '验证失败', 'verify your account', '点击此处', 'click here',
        '您的账户', '安全警告', 'security alert', '中奖', 'winner',
        '汇款', '转账', 'wire transfer', '发票', 'invoice',
        '退款', 'refund', '付款确认', 'payment confirmation',
        '报价', 'quotation', '合同', 'contract',
    ]
    MED_RISK = [
        '通知', 'notification', '更新', 'update', '确认', 'confirm',
        '重要', 'important', '提醒', 'reminder',
    ]

    subject = email_data.get('subject', '').lower()
    for kw in HIGH_RISK:
        if kw.lower() in subject:
            result['matched_keywords'].append(kw)
            result['risk_score'] += 2.0
    for kw in MED_RISK:
        if kw.lower() in subject:
            result['matched_keywords'].append(kw)
            result['risk_score'] += 0.5

    result['risk_level'] = risk_label(result['risk_score'])
    if result['matched_keywords']:
        result['warnings'].append(
            f"主题含高风险关键词: {', '.join(set(result['matched_keywords']))}"
        )
    return result


def detect_suspicious_attachments(email_data: Dict[str, Any]) -> Dict[str, Any]:
    """
    增强附件风险检测：
    - 双扩展名（report.pdf.exe）
    - 伪装为图片/文档的可执行文件
    - 压缩包内含可执行文件
    - 宏文档（.docm/.xlsm）
    """
    result = {
        'suspicious': [],
        'risk_level': 'low',
        'risk_score': 0.0,
        'warnings': []
    }

    EXEC_EXT = {
        '.exe', '.dll', '.bat', '.cmd', '.msi', '.vbs', '.vbe',
        '.js', '.jse', '.ps1', '.psm1', '.com', '.scr', '.hta',
        '.pif', '.lnk', '.jar', '.wsf', '.wsh', '.reg'
    }
    MACRO_EXT = {'.docm', '.xlsm', '.pptm', '.dotm', '.xltm', '.potm'}
    DOUBLE_EXT_RE = re.compile(
        r'\.(pdf|doc|docx|xls|xlsx|jpg|png|txt|zip)\.'
        r'(exe|bat|cmd|vbs|js|ps1|scr|com|pif|lnk)$',
        re.IGNORECASE
    )

    for att in email_data.get('attachments', []):
        fname = att.get('filename', '')
        ext   = att.get('extension', '').lower()
        issues = []

        # 双扩展名
        if DOUBLE_EXT_RE.search(fname):
            issues.append(f'双扩展名伪装: {fname}')
            result['risk_score'] += 5.0

        # 可执行文件
        if ext in EXEC_EXT:
            issues.append(f'可执行文件: {fname}')
            result['risk_score'] += 4.0

        # 宏文档
        if ext in MACRO_EXT:
            issues.append(f'含宏的文档（可能自动执行恶意代码）: {fname}')
            result['risk_score'] += 3.0

        # 压缩包内含可执行文件
        archive_contents = att.get('archive_contents', [])
        exec_in_archive = [
            f for f in archive_contents
            if os.path.splitext(f)[1].lower() in EXEC_EXT
        ]
        if exec_in_archive:
            issues.append(
                f'压缩包 {fname} 内含可执行文件: {", ".join(exec_in_archive[:3])}'
            )
            result['risk_score'] += 4.0

        if issues:
            result['suspicious'].append({'filename': fname, 'issues': issues})
            result['warnings'].extend(issues)

    result['risk_level'] = risk_label(result['risk_score'])
    return result


def detect_time_anomaly(email_data: Dict[str, Any]) -> Dict[str, Any]:
    """
    检测邮件头时间异常：
    - Date 与 Received 时间戳相差过大（>24h，说明被延迟/伪造）
    - Date 为未来时间
    """
    result = {
        'risk_level': 'low',
        'risk_score': 0.0,
        'warnings': []
    }

    date_str = email_data.get('date', '')
    if not date_str:
        return result

    try:
        from email.utils import parsedate_to_datetime
        mail_dt = parsedate_to_datetime(date_str)
        now = datetime.now(timezone.utc)

        # 未来时间
        if mail_dt > now:
            diff_h = (mail_dt - now).total_seconds() / 3600
            result['risk_score'] += 3.0
            result['warnings'].append(
                f'邮件日期({date_str})早于当前时间 {diff_h:.1f} 小时，疑似伪造时间戳'
            )

        # 与 Received 链中首个时间戳对比
        headers = email_data.get('headers', {})
        received_list = headers.get('received', [])
        if isinstance(received_list, str):
            received_list = [received_list]

        for rcv in received_list[:1]:
            # Received 头尾部通常有时间戳（;后面）
            ts_match = re.search(r';\s*(.+)$', rcv)
            if ts_match:
                try:
                    rcv_dt = parsedate_to_datetime(ts_match.group(1).strip())
                    diff_sec = abs((mail_dt - rcv_dt).total_seconds())
                    if diff_sec > 86400:  # 超过24小时
                        result['risk_score'] += 2.0
                        result['warnings'].append(
                            f'邮件Date与Received时间戳相差 {diff_sec/3600:.1f} 小时，可能被延迟或时间伪造'
                        )
                except Exception:
                    pass
    except Exception:
        pass

    result['risk_level'] = risk_label(result['risk_score'])
    return result


def extract_urls(email_data: Dict[str, Any]) -> Dict[str, Any]:
    """
    从邮件中提取所有URL并进行安全分析
    """
    result = {
        'urls': {
            'text': [],      # 纯文本中的URL
            'html': [],      # HTML内容中的URL
            'attachments': [] # 附件内容中的URL
        },
        'suspicious_links': [],  # 可疑的超链接
        'risk_level': 'low',
        'risk_score': 0.0,
        'warnings': []
    }
    
    try:
        # URL正则表达式模式
        url_pattern = r'(?i)\b((?:https?://|www\d{0,3}[.]|[a-z0-9.\-]+[.][a-z]{2,4}/)(?:[^\s()<>]+|\(([^\s()<>]+|(\([^\s()<>]+\)))*\))+(?:\(([^\s()<>]+|(\([^\s()<>]+\)))*\)|[^\s`!()\[\]{};:\'".,<>?«»""'']))'
        
        def analyze_link_safety(display_text: str, actual_url: str, email_domains: List[str]) -> Dict[str, Any]:
            """分析超链接的安全性"""
            res = {
                'display_text': display_text,
                'actual_url': actual_url,
                'risk_level': 'low',
                'risk_score': 0.0,
                'reasons': []
            }

            display_domain = extract_url_domain(display_text) if re.search(url_pattern, display_text) else None
            actual_domain  = extract_url_domain(actual_url)

            if display_domain and actual_domain and display_domain != actual_domain:
                res['risk_score'] += 3.0
                res['reasons'].append(
                    f"超链接显示域名({display_domain})与实际域名({actual_domain})不匹配"
                )

            for email_domain in email_domains:
                if actual_domain:
                    analysis = analyze_domain_similarity(actual_domain, email_domain)
                    if analysis['risk'] == 'high':
                        res['risk_score'] += 2.5
                        res['reasons'].append(
                            f"URL域名({actual_domain})与邮件域名({email_domain})高度相似: "
                            f"{analysis['relationship_type']}"
                        )

            try:
                from urllib.parse import urlparse, parse_qs
                parsed = urlparse(actual_url)
                if parsed.port and parsed.port not in (80, 443):
                    res['risk_score'] += 2.0
                    res['reasons'].append(f"使用非标准端口: {parsed.port}")
                if actual_url.count('%') > 5:
                    res['risk_score'] += 1.5
                    res['reasons'].append("URL过度编码，可能试图隐藏真实地址")
                redirect_params = {'url', 'redirect', 'goto', 'link', 'return', 'target'}
                found_redirects = set(parse_qs(parsed.query).keys()) & redirect_params
                if found_redirects:
                    res['risk_score'] += 2.0
                    res['reasons'].append(f"包含重定向参数: {', '.join(found_redirects)}")
            except Exception as e:
                print(f"URL分析失败: {str(e)}")

            if res['risk_score'] >= 3.0:
                res['risk_level'] = 'high'
            elif res['risk_score'] >= 1.5:
                res['risk_level'] = 'medium'
            return res
        
        # 获取邮件相关的域名列表
        email_domains = set()
        for sender in email_data['from']:
            domain_match = re.search(r'@([\w.-]+)', sender)
            if domain_match:
                email_domains.add(domain_match.group(1).lower())
        
        # 从纯文本中提取URL
        if email_data['body_text']:
            text_urls = re.findall(url_pattern, email_data['body_text'])
            result['urls']['text'] = [url[0] for url in text_urls]
        
        # 从HTML内容中提取URL和分析超链接
        if email_data['body_html']:
            try:
                soup = BeautifulSoup(email_data['body_html'], 'html.parser')
                
                # 分析所有超链接
                for link in soup.find_all('a'):
                    href = link.get('href')
                    if href and not href.startswith('mailto:'):
                        display_text = link.get_text(strip=True)
                        analysis = analyze_link_safety(display_text, href, list(email_domains))
                        
                        result['urls']['html'].append(href)
                        if analysis['risk_level'] != 'low':
                            result['suspicious_links'].append(analysis)
                            result['risk_score'] += analysis['risk_score']
                
                # 提取图片链接
                for img in soup.find_all('img'):
                    src = img.get('src')
                    if src and not src.startswith('data:'):
                        result['urls']['html'].append(src)
                
            except Exception as e:
                print(f"解析HTML内容时出错: {str(e)}")
        
        # 从附件中提取URL
        for attachment in email_data['attachments']:
            if attachment.get('text_preview'):
                attachment_urls = re.findall(url_pattern, attachment['text_preview'])
                result['urls']['attachments'].extend([url[0] for url in attachment_urls])
        
        # 去重并清理URL
        for source in result['urls']:
            result['urls'][source] = list(set(result['urls'][source]))
            result['urls'][source] = [url.rstrip('.,;:\'\"!?') for url in result['urls'][source]]
            result['urls'][source] = [url for url in result['urls'][source] if url]
        
        # 设置最终风险等级
        if result['risk_score'] >= 5.0:
            result['risk_level'] = 'high'
            result['warnings'].append("发现多个高风险URL，疑似钓鱼邮件！")
        elif result['risk_score'] >= 2.5:
            result['risk_level'] = 'medium'
            result['warnings'].append("发现可疑URL，请谨慎点击")
        
        # 添加统计信息
        total_urls = len(result['urls']['text']) + len(result['urls']['html']) + len(result['urls']['attachments'])
        result['warnings'].append(f"总计发现 {total_urls} 个URL，其中 {len(result['suspicious_links'])} 个可疑")
        
    except Exception as e:
        print(f"URL分析失败: {str(e)}")
        result['warnings'].append(f"URL分析过程出错: {str(e)}")
    
    return result

def detect_hidden_content(email_data: Dict[str, Any]) -> Dict[str, List[Dict[str, Any]]]:
    """
    检测邮件中的隐藏内容和跟踪器
    
    Args:
        email_data: 邮件解析数据
        
    Returns:
        包含检测结果的字典
    """
    findings = {
        'hidden_content': [],
        'tracking_elements': [],
        'suspicious_urls': [],
        'risk_level': 'low',
        'risk_score': 0.0,
        'warnings': []
    }
    
    try:
        if email_data['body_html']:
            # 尝试多种编码方式解析HTML内容
            html_content = email_data['body_html']
            if isinstance(html_content, bytes):
                encodings = ['utf-8', 'gbk', 'gb2312', 'gb18030', 'big5']
                for encoding in encodings:
                    try:
                        html_content = html_content.decode(encoding)
                        break
                    except:
                        continue
            
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # 1. 检查隐藏的图片和跟踪像素
            hidden_images = soup.find_all('img', style=lambda x: x and any(style in x.lower() for style in [
                'display:none', 'display: none',
                'visibility:hidden', 'visibility: hidden',
                'opacity:0', 'opacity: 0',
                'width:1px', 'width: 1px',
                'height:1px', 'height: 1px'
            ]))
            
            for img in hidden_images:
                src = img.get('src', '')
                findings['tracking_elements'].append({
                    'type': 'hidden_tracking_pixel',
                    'url': src,
                    'reason': '隐藏的跟踪图片',
                    'details': f'样式: {img.get("style", "")}'
                })
                findings['risk_score'] += 3.0
            
            # 2. 检查所有图片的尺寸属性
            all_images = soup.find_all('img')
            for img in all_images:
                width = img.get('width', '').strip()
                height = img.get('height', '').strip()
                if (width == '1' and height == '1') or (width == '0' and height == '0'):
                    src = img.get('src', '')
                    findings['tracking_elements'].append({
                        'type': 'tracking_pixel',
                        'url': src,
                        'reason': f'{width}x{height}像素跟踪图片'
                    })
                    findings['risk_score'] += 2.5
            
            # 3. 检查隐藏的内容
            hidden_elements = soup.find_all(lambda tag: tag.get('style') and any(style in tag.get('style', '').lower() for style in [
                'display:none', 'display: none',
                'visibility:hidden', 'visibility: hidden',
                'opacity:0', 'opacity: 0',
                'font-size:0', 'font-size: 0'
            ]))
            
            for element in hidden_elements:
                content = element.get_text().strip()
                if content:
                    findings['hidden_content'].append({
                        'type': 'hidden_element',
                        'content': content,
                        'style': element.get('style', ''),
                        'reason': '使用CSS隐藏的内容'
                    })
                    findings['risk_score'] += 2.0
            
            # 4. 检查可疑的URL和链接
            links = soup.find_all(['a', 'img', 'form'])
            for link in links:
                url = link.get('href') or link.get('src') or link.get('action', '')
                if url:
                    try:
                        from urllib.parse import urlparse, parse_qs
                        parsed = urlparse(url)
                        
                        # 检查非标准端口
                        if parsed.port and parsed.port not in (80, 443):
                            findings['suspicious_urls'].append({
                                'type': 'suspicious_port',
                                'url': url,
                                'port': parsed.port,
                                'reason': f'使用非标准端口 {parsed.port}'
                            })
                            findings['risk_score'] += 2.5
                        
                        # 检查可疑参数
                        query_params = parse_qs(parsed.query)
                        tracking_params = {'uid', 'user', 'id', 'email', 'track', 'open', 'click'}
                        found_params = set(query_params.keys()) & tracking_params
                        if found_params:
                            findings['tracking_elements'].append({
                                'type': 'tracking_parameters',
                                'url': url,
                                'params': list(found_params),
                                'reason': '包含跟踪参数'
                            })
                            findings['risk_score'] += 1.5
                        
                        # 检查URL编码过度
                        if url.count('%') > 5:
                            findings['suspicious_urls'].append({
                                'type': 'heavily_encoded_url',
                                'url': url,
                                'reason': 'URL过度编码，可能试图隐藏真实地址'
                            })
                            findings['risk_score'] += 2.0
                        
                    except Exception as e:
                        print(f"URL分析失败: {str(e)}")
            
            # 5. 检查外部资源加载
            external_resources = soup.find_all(['script', 'iframe', 'img', 'link'])
            for resource in external_resources:
                src = resource.get('src') or resource.get('href', '')
                if src and ('track' in src.lower() or 'beacon' in src.lower() or 'pixel' in src.lower()):
                    findings['tracking_elements'].append({
                        'type': 'external_tracker',
                        'url': src,
                        'tag': resource.name,
                        'reason': '外部跟踪资源'
                    })
                    findings['risk_score'] += 2.0
            
            # 更新风险等级
            if findings['risk_score'] >= 5.0:
                findings['risk_level'] = 'high'
                findings['warnings'].append('发现多个高风险跟踪器和隐藏内容，疑似钓鱼邮件！')
            elif findings['risk_score'] >= 3.0:
                findings['risk_level'] = 'medium'
                findings['warnings'].append('发现可疑的跟踪器或隐藏内容')
            
            # 添加统计信息
            findings['warnings'].append(f"发现 {len(findings['tracking_elements'])} 个跟踪元素")
            findings['warnings'].append(f"发现 {len(findings['hidden_content'])} 个隐藏内容")
            findings['warnings'].append(f"发现 {len(findings['suspicious_urls'])} 个可疑URL")
        
    except Exception as e:
        print(f"检测隐藏内容失败: {str(e)}")
        findings['warnings'].append(f"内容检测过程出错: {str(e)}")
    
    return findings

def verify_email_auth(email_data: Dict[str, Any]) -> Dict[str, Any]:
    """
    验证邮件认证信息，检查SPF、DKIM和DMARC
    
    Args:
        email_data: 邮件解析数据
        
    Returns:
        包含认证结果的字典
    """
    auth_results = {
        'spf': {
            'status': 'unknown',  # pass, fail, softfail, neutral, unknown
            'domain': '',
            'ip': '',
            'explanation': ''
        },
        'dkim': {
            'status': 'unknown',  # pass, fail, neutral, unknown
            'domain': '',
            'selector': '',
            'explanation': ''
        },
        'dmarc': {
            'status': 'unknown',  # pass, fail, none, unknown
            'domain': '',
            'policy': '',
            'explanation': ''
        },
        'authentication_results': '',  # 原始认证结果头信息
        'risk_level': 'unknown',      # high, medium, low, unknown
        'risk_score': 0.0,
        'warnings': []
    }
    
    try:
        # 1. 获取认证相关的头信息
        headers = {
            'authentication_results': '',  # Authentication-Results
            'received_spf': '',           # Received-SPF
            'dkim_signature': '',         # DKIM-Signature
            'dkim_results': '',           # DKIM验证结果
            'arc_authentication_results': '', # ARC-Authentication-Results
            'received': []                # Received链
        }
        
        # 2. 解析SPF记录
        def parse_spf_result(spf_header: str) -> None:
            if 'pass' in spf_header.lower():
                auth_results['spf']['status'] = 'pass'
            elif 'fail' in spf_header.lower():
                auth_results['spf']['status'] = 'fail'
            elif 'softfail' in spf_header.lower():
                auth_results['spf']['status'] = 'softfail'
            elif 'neutral' in spf_header.lower():
                auth_results['spf']['status'] = 'neutral'
            
            # 提取发送域名和IP
            domain_match = re.search(r'domain=(\S+)', spf_header)
            if domain_match:
                auth_results['spf']['domain'] = domain_match.group(1)
            
            ip_match = re.search(r'client-ip=(\S+)', spf_header)
            if ip_match:
                auth_results['spf']['ip'] = ip_match.group(1)
        
        # 3. 解析DKIM签名
        def parse_dkim_result(dkim_header: str) -> None:
            if 'pass' in dkim_header.lower():
                auth_results['dkim']['status'] = 'pass'
            elif 'fail' in dkim_header.lower():
                auth_results['dkim']['status'] = 'fail'
            
            # 提取DKIM域名和选择器
            domain_match = re.search(r'd=([^;]+)', dkim_header)
            if domain_match:
                auth_results['dkim']['domain'] = domain_match.group(1)
            
            selector_match = re.search(r's=([^;]+)', dkim_header)
            if selector_match:
                auth_results['dkim']['selector'] = selector_match.group(1)
        
        # 4. 解析DMARC结果
        def parse_dmarc_result(auth_header: str) -> None:
            dmarc_match = re.search(r'dmarc=(\S+)', auth_header.lower())
            if dmarc_match:
                status = dmarc_match.group(1)
                auth_results['dmarc']['status'] = status
                
                # 提取DMARC策略
                policy_match = re.search(r'p=(\S+)', auth_header)
                if policy_match:
                    auth_results['dmarc']['policy'] = policy_match.group(1)
        
        # 5. 分析认证头信息
        headers = email_data.get('headers', {})

        # 读取 Authentication-Results（可能有多条）
        auth_header = headers.get('authentication-results', '')
        if auth_header:
            auth_results['authentication_results'] = auth_header
            parse_spf_result(auth_header)
            parse_dkim_result(auth_header)
            parse_dmarc_result(auth_header)

        # 读取 Received-SPF（备用 SPF 来源）
        received_spf = headers.get('received-spf', '')
        if received_spf and auth_results['spf']['status'] == 'unknown':
            parse_spf_result(received_spf)

        # 读取 DKIM-Signature
        dkim_sig = headers.get('dkim-signature', '')
        if dkim_sig and auth_results['dkim']['status'] == 'unknown':
            parse_dkim_result(dkim_sig)
        
        # 6. 评估风险等级
        risk_score = 0.0
        
        # SPF检查
        if auth_results['spf']['status'] == 'fail':
            risk_score += 3.0
            auth_results['warnings'].append('SPF验证失败，发件人域名可能被伪造')
        elif auth_results['spf']['status'] == 'softfail':
            risk_score += 1.5
            auth_results['warnings'].append('SPF软失败，发件人域名可疑')
        
        # DKIM检查
        if auth_results['dkim']['status'] == 'fail':
            risk_score += 3.0
            auth_results['warnings'].append('DKIM签名验证失败，邮件可能被篡改')
        elif auth_results['dkim']['status'] == 'unknown':
            risk_score += 1.0
            auth_results['warnings'].append('未找到DKIM签名')
        
        # DMARC检查
        if auth_results['dmarc']['status'] == 'fail':
            risk_score += 3.0
            auth_results['warnings'].append('DMARC验证失败，不符合发件人域名的邮件策略')
        elif auth_results['dmarc']['status'] == 'none':
            risk_score += 1.0
            auth_results['warnings'].append('发件人域名未配置DMARC策略')
        
        # 发件人域名检查
        sender_domain = ''
        if email_data['from']:
            sender_match = re.search(r'@([\w.-]+)', email_data['from'][0])
            if sender_match:
                sender_domain = sender_match.group(1)
                
                # 检查发件人域名与认证域名是否匹配
                if sender_domain != auth_results['spf']['domain']:
                    risk_score += 2.0
                    auth_results['warnings'].append('发件人域名与SPF认证域名不匹配')
                if sender_domain != auth_results['dkim']['domain']:
                    risk_score += 2.0
                    auth_results['warnings'].append('发件人域名与DKIM签名域名不匹配')
        
        # 设置最终风险等级
        auth_results['risk_score'] = risk_score
        if risk_score >= 5.0:
            auth_results['risk_level'] = 'high'
        elif risk_score >= 2.0:
            auth_results['risk_level'] = 'medium'
        elif risk_score > 0:
            auth_results['risk_level'] = 'low'
        
    except Exception as e:
        print(f"验证邮件认证信息失败: {str(e)}")
        auth_results['warnings'].append(f"认证验证过程出错: {str(e)}")
    
    return auth_results

def display_report(email_data: Dict[str, Any]) -> None:
    """显示邮件分析报告（彩色高亮 + 综合评分）"""

    # ── ANSI 颜色常量 ──
    R  = '\033[1;31m'   # 红（高危）
    Y  = '\033[1;33m'   # 黄（中危）
    G  = '\033[1;32m'   # 绿（低危/正常）
    C  = '\033[1;36m'   # 青（信息）
    B  = '\033[1;34m'   # 蓝（标题）
    W  = '\033[1;37m'   # 白粗体
    RS = '\033[0m'      # 重置

    def clr(text, level):
        """按风险等级着色"""
        m = {'critical': R, 'high': R, 'medium': Y, 'low': G, 'unknown': W}
        return f"{m.get(level.lower(), W)}{text}{RS}"

    def warn(msg, level='high'):
        c = R if level in ('critical', 'high') else Y
        print(f"  {c}⚠  {msg}{RS}")

    def ok(msg):
        print(f"  {G}✔  {msg}{RS}")

    def section(num, title):
        print(f"\n{B}{'─'*60}{RS}")
        print(f"{W}[{num}] {title}{RS}")

    # ══════════════════════════════════════════════
    # 先运行所有检测，收集分数，最后汇总
    # ══════════════════════════════════════════════
    auth_r     = verify_email_auth(email_data)
    domain_r   = check_similar_domains(email_data)
    spoof_r    = detect_spoofed_sender(email_data)
    reg_r      = analyze_domain_registration(email_data)
    hidden_r   = detect_hidden_content(email_data)
    url_r      = extract_urls(email_data)
    att_r      = detect_suspicious_attachments(email_data)
    subj_r     = detect_suspicious_subject(email_data)
    homo_r     = detect_homograph_attack(email_data)
    time_r     = detect_time_anomaly(email_data)


    # ══════════════════════════════════════════════
    # 综合评分（满分 100）
    # ──────────────────────────────────────────────
    # 权重设计依据（安全实践）：
    #
    # 【极高权重 - 单独触发即高度可疑】
    #   同形字攻击(15)：几乎100%恶意，无正常使用场景
    #   附件威胁(15)：可执行文件/宏/双扩展名，直接危害最高
    #
    # 【高权重 - 强信号】
    #   发件人伪造(15)：SPF/Received链不匹配，确认身份欺骗
    #   邮件认证(15)：SPF/DKIM/DMARC三重失败，发件人可信度极低
    #   域名仿冒(12)：字符替换/高相似域名，典型钓鱼手法
    #
    # 【中权重 - 辅助信号】
    #   URL风险(10)：链接显示与实际不符，辅助判断
    #   隐藏内容(8)：跟踪像素等，钓鱼邮件常用
    #   域名年龄(8)：新注册域名，单独价值中等
    #
    # 【低权重 - 弱信号，配合其他使用】
    #   主题关键词(5)：正常邮件也可能触发
    #   时间异常(5)：单独出现可能是服务器问题
    # ──────────────────────────────────────────────
    # 各维度格式：(原始分, 原始满分, 权重)
    WEIGHTS = {
        'homo':    (homo_r['risk_score'],    4.0,  15),  # 同形字：极高危
        'att':     (att_r['risk_score'],    10.0,  15),  # 附件威胁：极高危
        'spoof':   (spoof_r['risk_score'],   8.0,  15),  # 发件伪造：高危
        'auth':    (auth_r['risk_score'],    9.0,  15),  # 邮件认证：高危
        'domain':  (domain_r['risk_score'],  6.0,  12),  # 域名仿冒：高危
        'url':     (url_r['risk_score'],     8.0,  10),  # URL风险：中危
        'hidden':  (hidden_r['risk_score'],  8.0,   8),  # 隐藏内容：中危
        'reg':     (reg_r['risk_score'],     5.0,   8),  # 域名年龄：中危
        'subj':    (subj_r['risk_score'],    6.0,   5),  # 主题词：弱信号
        'time':    (time_r['risk_score'],    4.0,   5),  # 时间异常：弱信号
    }
    # 各权重之和 = 108，压缩到100

    total_score = 0.0
    for key, (score, max_raw, weight) in WEIGHTS.items():
        normalized = min(score / max_raw, 1.0) * weight if max_raw > 0 else 0
        total_score += normalized

    # ── 联动加分：多个强信号同时命中时额外加分 ──
    # 钓鱼组合1：域名仿冒 + 域名年龄短（典型新建仿冒域名）
    if domain_r['risk_level'] == 'high' and reg_r.get('sender_domain', {}).get('age_days', 999) < 90:
        bonus = 15.0
        total_score += bonus

    # 钓鱼组合2：发件伪造 + 认证失败（双重身份欺骗）
    if spoof_r['is_spoofed'] and auth_r['spf']['status'] in ('fail', 'softfail'):
        total_score += 10.0

    # 钓鱼组合3：域名仿冒 + 可疑URL（视觉欺骗配合链接劫持）
    if domain_r['risk_level'] == 'high' and url_r['risk_level'] in ('high', 'medium'):
        total_score += 8.0

    # 钓鱼组合4：同形字 + 任意其他高危信号（几乎确认恶意）
    if homo_r['risk_score'] >= 4 and any([
        auth_r['risk_level'] == 'high',
        domain_r['risk_level'] == 'high',
        spoof_r['is_spoofed'],
    ]):
        total_score += 12.0

    total_score = min(total_score, 100.0)

    # ── 风险等级阈值（相比旧版更严格）──
    if total_score >= 65:
        overall_level = 'critical'
    elif total_score >= 40:
        overall_level = 'high'
    elif total_score >= 20:
        overall_level = 'medium'
    else:
        overall_level = 'low'

    # ── 恶意判定：满足任一条件 ──
    # 1. 综合评分高
    # 2. 极强单一信号（同形字/可执行附件）
    # 3. 三个以上高危信号联合
    high_signals = sum([
        auth_r['spf']['status'] in ('fail',),
        auth_r['dkim']['status'] in ('fail',),
        domain_r['risk_level'] == 'high',
        spoof_r['is_spoofed'],
        att_r['risk_level'] in ('high', 'critical'),
        homo_r['risk_score'] >= 4,
        reg_r.get('sender_domain', {}).get('age_days', 999) < 30,
        url_r['risk_level'] == 'high',
    ])
    is_malicious = (
        total_score >= 55
        or homo_r['risk_score'] >= 4          # 同形字单独触发
        or att_r['risk_score'] >= 8            # 可执行附件/双扩展名单独触发
        or high_signals >= 4                   # 4个及以上高危信号
        or (domain_r['risk_level'] == 'high'   # 域名仿冒 + 新域名（≤30天）
            and reg_r.get('sender_domain', {}).get('age_days', 999) <= 30)
    )


    # ══════════════════════════════════════════════
    # 顶部横幅
    # ══════════════════════════════════════════════
    print(f"\n{B}{'═'*60}{RS}")
    print(f"{W}  📧  邮件安全分析报告{RS}")
    print(f"{B}{'═'*60}{RS}")

    if is_malicious:
        print(f"\n{R}{'█'*60}{RS}")
        print(f"{R}  ‼‼  高度疑似恶意邮件！请勿点击任何链接或附件！  ‼‼{RS}")
        print(f"{R}{'█'*60}{RS}")
    
    # 综合评分展示
    bar_filled = int(total_score / 5)
    bar = f"{'█' * bar_filled}{'░' * (20 - bar_filled)}"
    score_color = R if total_score >= 60 else (Y if total_score >= 30 else G)
    print(f"\n  综合风险评分: {score_color}{total_score:.1f} / 100  [{bar}]{RS}")
    print(f"  综合风险等级: {clr(overall_level.upper(), overall_level)}")

    # ══════════════════════════════════════════════
    # [1] 基本信息
    # ══════════════════════════════════════════════
    section(1, '基本信息')
    print(f"  发件人: {W}{', '.join(email_data['from']) or '未知'}{RS}")
    print(f"  收件人: {', '.join(email_data['to']) or '未知'}")
    if email_data['cc']:
        print(f"  抄送:   {', '.join(email_data['cc'])}")
    if email_data['reply_to']:
        from_d  = extract_email_domain(email_data['from'][0]) if email_data['from'] else ''
        reply_d = extract_email_domain(email_data['reply_to'][0])
        rt_str  = ', '.join(email_data['reply_to'])
        if from_d and reply_d and from_d != reply_d:
            print(f"  回复地址: {R}{rt_str}  ← 与发件人域名不同！{RS}")
        else:
            print(f"  回复地址: {rt_str}")
    print(f"  主题:   {W}{email_data['subject'] or '未知'}{RS}")
    print(f"  日期:   {email_data['date'] or '未知'}")

    # 时间异常
    if time_r['warnings']:
        for w in time_r['warnings']:
            warn(w, 'medium')

    # 主题高危关键词
    if subj_r['matched_keywords']:
        warn(f"主题含高风险关键词: {', '.join(set(subj_r['matched_keywords']))}", 'medium')

    # 同形字攻击
    if homo_r['suspicious']:
        for w in homo_r['warnings']:
            warn(w, 'high')

    # ══════════════════════════════════════════════
    # [2] 邮件认证（SPF / DKIM / DMARC）
    # ══════════════════════════════════════════════
    section(2, f'邮件认证  风险: {clr(auth_r["risk_level"].upper(), auth_r["risk_level"])}  得分贡献: {min(auth_r["risk_score"]/9,1)*20:.1f}/20')

    def auth_status(name, status):
        if status == 'pass':
            ok(f"{name}: {G}PASS{RS}")
        elif status in ('fail', 'softfail'):
            warn(f"{name}: {status.upper()} — 验证失败，发件人可能被伪造", 'high')
        elif status == 'unknown':
            print(f"  {Y}?  {name}: 未知/未找到{RS}")
        else:
            print(f"  {W}·  {name}: {status}{RS}")

    auth_status('SPF',   auth_r['spf']['status'])
    auth_status('DKIM',  auth_r['dkim']['status'])
    auth_status('DMARC', auth_r['dmarc']['status'])

    for w in auth_r['warnings']:
        warn(w)

    # ══════════════════════════════════════════════
    # [3] 发件人真实性
    # ══════════════════════════════════════════════
    section(3, f'发件人真实性  风险: {clr(spoof_r["risk_level"].upper(), spoof_r["risk_level"])}  得分贡献: {min(spoof_r["risk_score"]/8,1)*15:.1f}/15')

    if spoof_r['is_spoofed']:
        for w in spoof_r['warnings']:
            warn(w)
    else:
        ok('未发现发件人伪造迹象')

    # ══════════════════════════════════════════════
    # [4] 域名相似度（仿冒检测）
    # ══════════════════════════════════════════════
    section(4, f'域名仿冒检测  风险: {clr(domain_r["risk_level"].upper(), domain_r["risk_level"])}  得分贡献: {min(domain_r["risk_score"]/6,1)*20:.1f}/20')

    if domain_r['similar_domains']:
        for item in domain_r['similar_domains']:
            warn(
                f"可疑域名对: {item['domain1']}  ↔  {item['domain2']}"
                f"  ({item['relationship']}  相似度:{item['similarity']:.0%})",
                'high'
            )
        for w in domain_r['warnings']:
            if '\n' in w:
                for line in w.split('\n'):
                    if line.strip():
                        warn(line.strip())
    else:
        ok('未发现可疑相似域名')

    # ══════════════════════════════════════════════
    # [5] 域名注册信息
    # ══════════════════════════════════════════════
    section(5, f'域名注册信息  风险: {clr(reg_r["risk_level"].upper(), reg_r["risk_level"])}  得分贡献: {min(reg_r["risk_score"]/5,1)*10:.1f}/10')

    for key, label in [('sender_domain', '发件人'), ('recipient_domain', '收件人')]:
        info = reg_r.get(key)
        if not info:
            continue
        age  = info.get('age_days')
        lvl  = info.get('risk_level', 'low')
        age_str = f"{age} 天" if age is not None else '未知'
        line = f"{label}域名 {info['domain']}  注册年龄: {age_str}"
        if lvl in ('critical', 'high'):
            warn(line, lvl)
        else:
            print(f"  {G}·{RS}  {line}  ({clr(lvl.upper(), lvl)})")

    for w in reg_r['warnings']:
        warn(w, 'medium')

    # ══════════════════════════════════════════════
    # [6] 隐藏内容 & 跟踪器
    # ══════════════════════════════════════════════
    section(6, f'隐藏内容/跟踪器  风险: {clr(hidden_r["risk_level"].upper(), hidden_r["risk_level"])}  得分贡献: {min(hidden_r["risk_score"]/8,1)*10:.1f}/10')

    if hidden_r['risk_level'] == 'low' and not hidden_r['tracking_elements']:
        ok('未发现隐藏内容或跟踪器')
    else:
        for item in hidden_r['tracking_elements']:
            warn(f"跟踪元素 [{item['type']}]: {item.get('url','')[:80]}  — {item['reason']}")
        for item in hidden_r['hidden_content']:
            warn(f"隐藏内容 [{item['type']}]: {item['content'][:60]}...", 'medium')
        for item in hidden_r['suspicious_urls']:
            warn(f"可疑URL [{item['type']}]: {item['url'][:80]}  — {item['reason']}", 'medium')

    # ══════════════════════════════════════════════
    # [7] URL 分析
    # ══════════════════════════════════════════════
    section(7, f'URL分析  风险: {clr(url_r["risk_level"].upper(), url_r["risk_level"])}  得分贡献: {min(url_r["risk_score"]/8,1)*10:.1f}/10')

    total_urls = sum(len(v) for v in url_r['urls'].values())
    print(f"  共发现 {total_urls} 个URL（文本:{len(url_r['urls']['text'])}  HTML:{len(url_r['urls']['html'])}  附件:{len(url_r['urls']['attachments'])}）")

    if url_r['suspicious_links']:
        print(f"  {R}发现 {len(url_r['suspicious_links'])} 个可疑链接:{RS}")
        for lnk in url_r['suspicious_links']:
            warn(f"显示: {lnk['display_text'][:40]}  →  实际: {lnk['actual_url'][:70]}")
            for reason in lnk['reasons']:
                print(f"       {Y}· {reason}{RS}")
    else:
        ok('未发现可疑链接')

    # ══════════════════════════════════════════════
    # [8] 附件分析
    # ══════════════════════════════════════════════
    section(8, f'附件分析  风险: {clr(att_r["risk_level"].upper(), att_r["risk_level"])}  得分贡献: {min(att_r["risk_score"]/10,1)*10:.1f}/10')

    if not email_data['attachments']:
        print('  无附件')
    else:
        for i, att in enumerate(email_data['attachments'], 1):
            fname = att.get('filename', '未知')
            size  = att.get('size', 0)
            ext   = att.get('extension', '')
            is_exec = att.get('is_executable', False)
            print(f"  {i}. {W}{fname}{RS}  ({size} 字节  {ext})")
            if att.get('hash_md5'):
                print(f"     MD5: {att['hash_md5']}")
            if is_exec:
                warn(f"可执行文件！请勿打开: {fname}", 'high')
            if att.get('is_archive') and att.get('archive_contents'):
                print(f"     压缩包内容: {', '.join(att['archive_contents'][:5])}")

        # 附件高危汇总
        for item in att_r.get('suspicious', []):
            for issue in item['issues']:
                warn(issue, 'high')

    # ══════════════════════════════════════════════
    # [9] 正文预览
    # ══════════════════════════════════════════════
    section(9, '正文预览')

    if email_data['body_text']:
        snippet = email_data['body_text'][:400]
        print(f"  {snippet}{'...' if len(email_data['body_text'])>400 else ''}")
        print(f"  {C}（纯文本共 {len(email_data['body_text'])} 字符）{RS}")
    elif email_data['body_html']:
        try:
            soup = BeautifulSoup(email_data['body_html'], 'html.parser')
            for tag in soup(['script', 'style']):
                tag.decompose()
            clean = soup.get_text(separator='\n', strip=True)
            snippet = clean[:400]
            print(f"  {snippet}{'...' if len(clean)>400 else ''}")
            print(f"  {C}（HTML转文本共 {len(clean)} 字符）{RS}")
        except Exception:
            print('  （HTML正文解析失败）')
    else:
        print('  未找到邮件正文')

    # ══════════════════════════════════════════════
    # 底部综合结论
    # ══════════════════════════════════════════════
    print(f"\n{B}{'═'*60}{RS}")
    print(f"{W}  综合评分: {score_color}{total_score:.1f}/100{RS}  {W}风险等级: {clr(overall_level.upper(), overall_level)}{RS}")

    # 各维度得分小结
    dim_names = {
        'homo':   '同形字攻击',
        'att':    '附件威胁',
        'spoof':  '发件伪造',
        'auth':   '邮件认证',
        'domain': '域名仿冒',
        'url':    'URL风险',
        'hidden': '隐藏内容',
        'reg':    '域名年龄',
        'subj':   '主题关键词',
        'time':   '时间异常',
    }
    print(f"\n  {'维度':<10} {'原始分':>6}  {'贡献分':>6}")
    print(f"  {'─'*28}")
    for key, (score, max_raw, weight) in WEIGHTS.items():
        contrib = min(score / max_raw, 1.0) * weight if max_raw > 0 else 0
        bar_c = R if contrib / weight >= 0.6 else (Y if contrib / weight >= 0.3 else G)
        print(f"  {dim_names[key]:<10} {score:>6.1f}  {bar_c}{contrib:>5.1f}{RS}/{weight}")

    if is_malicious:
        print(f"\n{R}{'█'*60}{RS}")
        print(f"{R}  ‼  结论：高度疑似恶意邮件，建议隔离并上报安全团队  ‼{RS}")
        print(f"{R}{'█'*60}{RS}")
    elif overall_level == 'medium':
        print(f"\n{Y}  ⚠  结论：存在可疑特征，请谨慎处理，勿轻易点击链接或附件{RS}")
    else:
        print(f"\n{G}  ✔  结论：未发现明显恶意特征，但仍需保持警惕{RS}")

    print(f"{B}{'═'*60}{RS}\n")


def check_similar_domains(email_data: Dict[str, Any]) -> Dict[str, Any]:
    """
    检查邮件中所有地址的域名相似度
    """
    result = {
        'similar_domains': [],
        'risk_level': 'low',
        'risk_score': 0.0,
        'warnings': []
    }
    
    try:
        # 收集所有域名
        domains = {
            'from': set(),
            'to': set(),
            'original_from': set()
        }

        for sender in email_data['from']:
            d = extract_email_domain(sender)
            if d:
                domains['from'].add(d)

        for recipient in email_data['to']:
            d = extract_email_domain(recipient)
            if d:
                domains['to'].add(d)

        if email_data['thread_info']['original_sender']:
            d = extract_email_domain(email_data['thread_info']['original_sender'])
            if d:
                domains['original_from'].add(d)

        # 比较所有域名组合
        # 1. 发件人域名与收件人域名比较
        for from_domain in domains['from']:
            for to_domain in domains['to']:
                analysis = analyze_domain_similarity(from_domain, to_domain)
                if analysis['risk'] == 'high':
                    result['similar_domains'].append({
                        'domain1': from_domain,
                        'domain2': to_domain,
                        'type': '发件人-收件人',
                        'relationship': analysis['relationship_type'],
                        'similarity': analysis['similarity']
                    })
                    result['risk_score'] += 3.0
                    warning = f"发现可疑域名组合: {from_domain} <-> {to_domain} ({analysis['relationship_type']})"
                    if analysis['relationship_type'] == '可疑的域名包含':
                        warning += "\n这是典型的钓鱼邮件手法，使用目标公司域名构造相似域名"
                    result['warnings'].append(warning)

        # 2. 发件人域名与原始发件人域名比较
        for from_domain in domains['from']:
            for orig_domain in domains['original_from']:
                analysis = analyze_domain_similarity(from_domain, orig_domain)
                if analysis['risk'] == 'high':
                    result['similar_domains'].append({
                        'domain1': from_domain,
                        'domain2': orig_domain,
                        'type': '发件人-原始发件人',
                        'relationship': analysis['relationship_type'],
                        'similarity': analysis['similarity']
                    })
                    result['risk_score'] += 3.0
        
        # 设置最终风险等级
        if result['risk_score'] >= 3.0:
            result['risk_level'] = 'high'
            if not result['warnings']:
                result['warnings'].append('发现高度相似的可疑域名，疑似钓鱼邮件！')
        elif result['risk_score'] >= 1.0:
            result['risk_level'] = 'medium'
            if not result['warnings']:
                result['warnings'].append('发现相似域名，建议仔细核实发件人身份')
        
    except Exception as e:
        print(f"域名相似度检查失败: {str(e)}")
    
    return result

def detect_spoofed_sender(email_data: Dict[str, Any]) -> Dict[str, Any]:
    """
    检测伪造的发件人
    
    Args:
        email_data: 邮件解析数据
        
    Returns:
        包含检测结果的字典
    """
    result = {
        'is_spoofed': False,
        'evidence': [],
        'risk_level': 'low',
        'risk_score': 0.0,
        'original_sender': '',
        'spoofed_sender': '',
        'warnings': []
    }
    
    try:
        # 1. 检查SPF验证结果
        auth_results = verify_email_auth(email_data)
        if auth_results['spf']['status'] == 'fail':
            result['is_spoofed'] = True
            result['evidence'].append('SPF验证失败')
            result['risk_score'] += 3.0
        
        # 2. 分析 Received 链
        received_chain = []
        headers = email_data.get('headers', {})
        received_val = headers.get('received', [])
        if isinstance(received_val, list):
            received_chain = received_val
        elif isinstance(received_val, str):
            received_chain = [received_val]

        if received_chain:
            last_received = received_chain[-1]
            claimed_domain = ''
            for sender in email_data['from']:
                domain_match = re.search(r'@([\w.-]+)', sender)
                if domain_match:
                    claimed_domain = domain_match.group(1).lower()
                    break

            if claimed_domain and claimed_domain not in last_received.lower():
                result['is_spoofed'] = True
                result['evidence'].append(
                    f'发件人声称来自 {claimed_domain}，但实际发送服务器不匹配'
                )
                result['risk_score'] += 2.5

        # 3. 检查 X-Fangmail-Spf 头
        if headers.get('x-fangmail-spf', '').lower() == 'fail':
            result['is_spoofed'] = True
            result['evidence'].append('防垃圾邮件系统SPF检查失败')
            result['risk_score'] += 2.0
        
        # 4. 检查 Reply-To 与 From 域名是否一致（Reply-To 劫持）
        if email_data['from'] and email_data['reply_to']:
            from_domain_match = re.search(r'@([\w.-]+)', email_data['from'][0])
            reply_domain_match = re.search(r'@([\w.-]+)', email_data['reply_to'][0])
            if from_domain_match and reply_domain_match:
                from_d = from_domain_match.group(1).lower()
                reply_d = reply_domain_match.group(1).lower()
                if from_d != reply_d:
                    result['is_spoofed'] = True
                    result['evidence'].append(
                        f'Reply-To域名({reply_d})与From域名({from_d})不一致，存在回复劫持风险'
                    )
                    result['risk_score'] += 2.5

        # 记录发件人
        if email_data['from']:
            result['spoofed_sender'] = email_data['from'][0]
        
        # 设置风险等级
        if result['risk_score'] >= 5.0:
            result['risk_level'] = 'high'
            result['warnings'].append('发现多个证据表明发件人身份被伪造，这很可能是一封钓鱼邮件！')
        elif result['risk_score'] >= 2.5:
            result['risk_level'] = 'medium'
            result['warnings'].append('发现可疑迹象表明发件人身份可能被伪造')
        
        # 添加详细分析
        if result['is_spoofed']:
            result['warnings'].extend([
                f"伪造的发件人: {result['spoofed_sender']}",
                "发现的证据:",
                *[f"- {evidence}" for evidence in result['evidence']]
            ])
            
    except Exception as e:
        print(f"检测伪造发件人失败: {str(e)}")
        result['warnings'].append(f"发件人检测过程出错: {str(e)}")
    
    return result

def check_domain_registration(domain: str) -> Dict[str, Any]:
    """
    检查域名的注册信息
    """
    result = {
        'domain': domain,
        'creation_date': None,
        'expiration_date': None,
        'last_updated': None,
        'age_days': None,
        'risk_level': 'low',
        'risk_score': 0.0,
        'warnings': []
    }
    
    try:
        # 预处理域名
        if not domain or '.' not in domain:
            raise ValueError(f"无效的域名格式: {domain}")
            
        # 移除可能的端口号和路径
        domain = domain.split(':')[0].split('/')[0]
        
        # 尝试 WHOIS 查询
        try:
            w = whois.whois(domain)
        except Exception as e:
            raise Exception(f"WHOIS查询失败: {str(e)}")
        
        if not w or not w.domain_name:
            raise Exception(f"无法获取域名信息")
        
        def ensure_datetime_with_timezone(dt):
            """确保日期时间对象带有时区信息"""
            if dt is None:
                return None
            if isinstance(dt, str):
                try:
                    # 尝试解析字符串为datetime对象
                    dt = datetime.strptime(dt, "%Y-%m-%d %H:%M:%S")
                except:
                    try:
                        dt = datetime.strptime(dt, "%Y-%m-%d")
                    except:
                        return None
            if not isinstance(dt, datetime):
                return None
            if dt.tzinfo is None:
                # 如果没有时区信息，假定为UTC
                dt = dt.replace(tzinfo=timezone.utc)
            return dt
        
        # 获取当前时间（带时区信息）
        now = datetime.now(timezone.utc)
        
        # 处理创建日期
        if w.creation_date:
            try:
                if isinstance(w.creation_date, list):
                    # 处理多个创建日期的情况，选择最早的有效日期
                    valid_dates = []
                    for date in w.creation_date:
                        dt = ensure_datetime_with_timezone(date)
                        if dt:
                            valid_dates.append(dt)
                    creation_date = min(valid_dates) if valid_dates else None
                else:
                    creation_date = ensure_datetime_with_timezone(w.creation_date)
                
                if creation_date:
                    result['creation_date'] = creation_date
                    age_days = (now - creation_date).days
                    result['age_days'] = age_days
                    
                    # 评估风险
                    if age_days <= 7:  # 一周内注册
                        result['risk_level'] = 'critical'
                        result['risk_score'] = 5.0
                        result['warnings'].append(f"域名 {domain} 注册时间极短（{age_days}天），非常可疑！")
                    elif age_days <= 30:  # 一个月内注册
                        result['risk_level'] = 'high'
                        result['risk_score'] = 4.0
                        result['warnings'].append(f"域名 {domain} 注册时间很短（{age_days}天），高度可疑")
                    elif age_days <= 90:  # 三个月内注册
                        result['risk_level'] = 'medium'
                        result['risk_score'] = 3.0
                        result['warnings'].append(f"域名 {domain} 注册时间较短（{age_days}天），需要注意")
                    elif age_days <= 365:  # 一年内注册
                        result['risk_level'] = 'low'
                        result['risk_score'] = 2.0
                        result['warnings'].append(f"域名 {domain} 注册时间不足一年（{age_days}天）")
            except Exception as e:
                result['warnings'].append(f"处理创建日期时出错: {str(e)}")
        
        # 处理过期日期
        if w.expiration_date:
            try:
                if isinstance(w.expiration_date, list):
                    # 处理多个过期日期的情况，选择最晚的有效日期
                    valid_dates = []
                    for date in w.expiration_date:
                        dt = ensure_datetime_with_timezone(date)
                        if dt:
                            valid_dates.append(dt)
                    expiration_date = max(valid_dates) if valid_dates else None
                else:
                    expiration_date = ensure_datetime_with_timezone(w.expiration_date)
                
                if expiration_date:
                    result['expiration_date'] = expiration_date
            except Exception as e:
                result['warnings'].append(f"处理过期日期时出错: {str(e)}")
        
        # 处理更新日期
        if w.updated_date:
            try:
                if isinstance(w.updated_date, list):
                    # 处理多个更新日期的情况，选择最近的有效日期
                    valid_dates = []
                    for date in w.updated_date:
                        dt = ensure_datetime_with_timezone(date)
                        if dt:
                            valid_dates.append(dt)
                    updated_date = max(valid_dates) if valid_dates else None
                else:
                    updated_date = ensure_datetime_with_timezone(w.updated_date)
                
                if updated_date:
                    result['last_updated'] = updated_date
            except Exception as e:
                result['warnings'].append(f"处理更新日期时出错: {str(e)}")
                
    except Exception as e:
        result['warnings'].append(f"无法获取域名 {domain} 的注册信息: {str(e)}")
        result['risk_level'] = 'unknown'
        result['risk_score'] = 0.0
    
    return result

def analyze_domain_registration(email_data: Dict[str, Any]) -> Dict[str, Any]:
    """
    分析邮件中涉及的域名注册信息
    """
    result = {
        'sender_domain': None,
        'recipient_domain': None,
        'domain_age_comparison': None,
        'risk_level': 'low',
        'risk_score': 0.0,
        'warnings': []
    }
    
    try:
        # 提取发件人域名
        sender_domain = None
        if email_data['from']:
            domain_match = re.search(r'@([\w.-]+)', email_data['from'][0])
            if domain_match:
                sender_domain = domain_match.group(1).lower()
                result['sender_domain'] = check_domain_registration_with_retry(sender_domain)
                if result['sender_domain'].get('risk_score'):
                    result['risk_score'] += result['sender_domain']['risk_score']
        
        # 提取收件人域名
        recipient_domain = None
        if email_data['to']:
            domain_match = re.search(r'@([\w.-]+)', email_data['to'][0])
            if domain_match:
                recipient_domain = domain_match.group(1).lower()
                result['recipient_domain'] = check_domain_registration_with_retry(recipient_domain)
        
        # 如果发件人和收件人域名不同，进行对比分析
        if (sender_domain and recipient_domain and 
            sender_domain != recipient_domain and 
            result['sender_domain'] and result['recipient_domain']):
            
            result['domain_age_comparison'] = {
                'age_difference_days': None,
                'warnings': []
            }
            
            sender_age = result['sender_domain'].get('age_days')
            recipient_age = result['recipient_domain'].get('age_days')
            
            if sender_age is not None and recipient_age is not None:
                try:
                    age_diff = recipient_age - sender_age
                    result['domain_age_comparison']['age_difference_days'] = age_diff
                    
                    if age_diff > 365:  # 收件人域名比发件人域名老一年以上
                        result['risk_score'] += 2.0
                        warning = (
                            f"发件人域名 {sender_domain} 比收件人域名 {recipient_domain} "
                            f"新 {age_diff} 天，这种情况在钓鱼邮件中很常见"
                        )
                        result['domain_age_comparison']['warnings'].append(warning)
                        result['warnings'].append(warning)
                except Exception as e:
                    result['warnings'].append(f"比较域名年龄时出错: {str(e)}")
        
        # 设置最终风险等级
        if result['risk_score'] >= 5.0:
            result['risk_level'] = 'critical'
        elif result['risk_score'] >= 4.0:
            result['risk_level'] = 'high'
        elif result['risk_score'] >= 3.0:
            result['risk_level'] = 'medium'
        
        # 合并所有警告信息
        if result['sender_domain'] and result['sender_domain'].get('warnings'):
            result['warnings'].extend(result['sender_domain']['warnings'])
        if result['recipient_domain'] and result['recipient_domain'].get('warnings'):
            result['warnings'].extend(result['recipient_domain']['warnings'])
            
    except Exception as e:
        result['warnings'].append(f"域名注册信息分析失败: {str(e)}")
    
    return result

def check_domain_registration_with_retry(domain: str, max_retries: int = 3) -> Dict[str, Any]:
    """带重试机制的域名注册信息检查"""
    for attempt in range(max_retries):
        try:
            return check_domain_registration(domain)
        except Exception as e:
            if attempt == max_retries - 1:
                return {
                    'domain': domain,
                    'creation_date': None,
                    'expiration_date': None,
                    'last_updated': None,
                    'age_days': None,
                    'risk_level': 'unknown',
                    'risk_score': 0.0,
                    'warnings': [f"多次尝试后仍无法获取域名信息: {str(e)}"]
                }
            print(f"第 {attempt + 1} 次尝试失败，准备重试...")
            time.sleep(2)  # 等待2秒后重试

def get_verified_path(path_type="any"):
    """获取已验证的指定类型路径"""
    while True:
        try:
            raw_input = input("路径输入（支持拖拽文件）: ").strip(' "\'')
            full_path = os.path.normpath(os.path.expanduser(raw_input))
            
            if not os.path.exists(full_path):
                raise FileNotFoundError("路径不存在")
                
            if path_type == "file" and not os.path.isfile(full_path):
                raise ValueError("必须选择文件")
                
            if path_type == "dir" and not os.path.isdir(full_path):
                raise ValueError("必须选择目录")
                
            return full_path
            
        except (FileNotFoundError, ValueError) as e:
            print(f"输入错误: {str(e)}")
        except Exception as e:
            print(f"未知错误: {str(e)}")


# 使用示例
if __name__ == "__main__":
    # test_file = r"/Users/admin/Downloads/1745443910003.eml"
    test_file = get_verified_path("file")

    try:
        print(f"开始解析邮件: {test_file}")
        print(f"文件大小: {os.path.getsize(test_file)} 字节")
        
        # 解析邮件
        email_data = parse_email(test_file)
        
        # 显示附件统计
        print(f"\n发现附件数量: {len(email_data['attachments'])}")
        
        # 显示报告
        display_report(email_data)
        
    except Exception as e:
        print(f"处理失败: {str(e)}")
        import traceback
        traceback.print_exc()
